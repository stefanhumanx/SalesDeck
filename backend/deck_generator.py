"""
Deck Generator - Core Logic

Given a PPTX template and a list of selected product slide indices, generates
a custom sponsorship deck PDF by:
1. Copying the template
2. Removing all non-selected product slides (indices 23-128)
3. Inserting a quote/summary slide after the front matter
4. Converting to PDF via LibreOffice
"""

import logging
import os
import shutil
import subprocess
import tempfile
from datetime import datetime
from typing import List, Dict, Any, Tuple
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# Slide ranges
FRONT_MATTER_END = 22  # slides 0-22 are static intro
PRODUCT_RANGE = range(23, 129)  # slides 23-128 are products
CLOSING_RANGE = range(129, 132)  # slides 129-131 are closing


def remove_slide(prs: Presentation, index: int) -> None:
    """
    Remove slide at 0-based index from Presentation object.
    
    Uses XML manipulation since python-pptx lacks a native remove_slide.
    
    Args:
        prs: Presentation object
        index: 0-based slide index to remove
    """
    try:
        xml_slides = prs.slides._sldIdLst
        if index >= len(xml_slides):
            logger.warning(f"Slide index {index} out of range")
            return
        
        sldId = xml_slides[index]
        rId = sldId.attrib.get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        
        if rId:
            prs.part.drop_rel(rId)
        
        xml_slides.remove(sldId)
        logger.debug(f"Removed slide at index {index}")
    except Exception as e:
        logger.error(f"Error removing slide at index {index}: {e}")
        raise


def add_quote_slide(
    prs: Presentation,
    sponsor_name: str,
    rep_name: str,
    products: List[Dict[str, Any]],
    insert_position: int
) -> None:
    """
    Insert a styled quote/sponsorship proposal slide.
    
    The slide includes:
    - Dark background
    - Title: "SPONSORSHIP PROPOSAL"
    - Sponsor name and rep name
    - Table with products and investments
    - Total investment row
    - Date generated
    
    Args:
        prs: Presentation object
        sponsor_name: Name of the sponsor
        rep_name: Name of the rep
        products: List of selected products with slide_index, name, price
        insert_position: Index where to insert the slide (0-based)
    """
    try:
        # Add a blank slide using the first available layout
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background to dark color (#1a1a2e)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(26, 26, 46)
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "SPONSORSHIP PROPOSAL"
        title_frame.word_wrap = True
        
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(54)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(245, 166, 35)  # Gold accent
        title_paragraph.alignment = PP_ALIGN.LEFT
        
        # Add sponsor info
        info_top = Inches(1.8)
        info_box = slide.shapes.add_textbox(left, info_top, width, Inches(0.8))
        info_frame = info_box.text_frame
        info_frame.clear()
        
        p = info_frame.paragraphs[0]
        p.text = f"Sponsor: {sponsor_name}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        
        p2 = info_frame.add_paragraph()
        p2.text = f"Representative: {rep_name}"
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(200, 200, 200)
        p2.alignment = PP_ALIGN.LEFT
        
        # Add products table
        table_top = Inches(2.8)
        rows = len(products) + 2  # +1 for header, +1 for total
        cols = 2
        
        left_table = Inches(1)
        width_table = Inches(8)
        height_table = Inches(0.4 * rows)
        
        table_shape = slide.shapes.add_table(rows, cols, left_table, table_top, width_table, height_table)
        table = table_shape.table
        
        # Set column widths
        table.columns[0].width = Inches(5)
        table.columns[1].width = Inches(3)
        
        # Header row
        header_cells = [table.cell(0, 0), table.cell(0, 1)]
        header_texts = ["Product", "Investment"]
        
        for cell, text in zip(header_cells, header_texts):
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 166, 35)  # Gold background
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(26, 26, 46)  # Dark text
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Product rows
        total_price = 0
        for idx, product in enumerate(products, start=1):
            name = product.get("name", "Unknown")
            price = product.get("price")
            
            # Product name cell
            cell = table.cell(idx, 0)
            cell.text = name
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(40, 40, 60)
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.LEFT
            
            # Price cell
            cell = table.cell(idx, 1)
            if price:
                cell.text = f"${price:,.0f}"
                total_price += price
            else:
                cell.text = "TBD"
            
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(40, 40, 60)
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.RIGHT
        
        # Total row
        last_idx = len(products) + 1
        total_name_cell = table.cell(last_idx, 0)
        total_name_cell.text = "Total Investment"
        total_name_cell.fill.solid()
        total_name_cell.fill.fore_color.rgb = RGBColor(245, 166, 35)
        
        p = total_name_cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 26, 46)
        p.alignment = PP_ALIGN.LEFT
        
        total_price_cell = table.cell(last_idx, 1)
        total_price_cell.text = f"${total_price:,.0f}"
        total_price_cell.fill.solid()
        total_price_cell.fill.fore_color.rgb = RGBColor(245, 166, 35)
        
        p = total_price_cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 26, 46)
        p.alignment = PP_ALIGN.RIGHT
        
        # Add date footer
        footer_top = Inches(6.5)
        footer_box = slide.shapes.add_textbox(left, footer_top, width, Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        footer_frame.word_wrap = False
        
        footer_paragraph = footer_frame.paragraphs[0]
        footer_paragraph.font.size = Pt(10)
        footer_paragraph.font.color.rgb = RGBColor(150, 150, 150)
        footer_paragraph.alignment = PP_ALIGN.RIGHT
        
        # Move the newly added slide to the correct position
        xml_slides = prs.slides._sldIdLst
        slides_count = len(xml_slides)
        
        if insert_position < slides_count:
            # Move the last added slide to insert_position
            new_slide_id = xml_slides[slides_count - 1]
            xml_slides.remove(new_slide_id)
            xml_slides.insert(insert_position, new_slide_id)
        
        logger.info(f"Added quote slide at position {insert_position}")
    
    except Exception as e:
        logger.error(f"Error adding quote slide: {e}")
        raise


def convert_to_pdf(pptx_path: str, output_dir: str) -> str:
    """
    Convert PPTX to PDF using LibreOffice headless.
    
    Tries multiple command variations in order:
    1. python /sessions/optimistic-peaceful-noether/mnt/.skills/skills/pptx/scripts/office/soffice.py
    2. soffice
    3. libreoffice
    
    Args:
        pptx_path: Path to the PPTX file
        output_dir: Directory to save the PDF
    
    Returns:
        Path to the generated PDF file
    
    Raises:
        RuntimeError: If conversion fails with all attempts
    """
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    commands = [
        [
            "python",
            "/sessions/optimistic-peaceful-noether/mnt/.skills/skills/pptx/scripts/office/soffice.py",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            pptx_path
        ],
        [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            pptx_path
        ],
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            pptx_path
        ]
    ]
    
    last_error = None
    
    for cmd in commands:
        try:
            logger.info(f"Attempting to convert PPTX to PDF with: {cmd[0]}")
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=120
            )
            
            if result.returncode == 0:
                # Find the generated PDF
                pdf_filename = Path(pptx_path).stem + ".pdf"
                pdf_path = os.path.join(output_dir, pdf_filename)
                
                if os.path.exists(pdf_path):
                    logger.info(f"Successfully converted PPTX to PDF: {pdf_path}")
                    return pdf_path
            else:
                last_error = f"Return code {result.returncode}: {result.stderr}"
                logger.debug(f"Command failed: {last_error}")
        
        except FileNotFoundError:
            last_error = f"Command not found: {cmd[0]}"
            logger.debug(last_error)
        except subprocess.TimeoutExpired:
            last_error = "Conversion timed out"
            logger.debug(last_error)
        except Exception as e:
            last_error = str(e)
            logger.debug(f"Error running command: {e}")
    
    raise RuntimeError(
        f"Failed to convert PPTX to PDF. Last error: {last_error}"
    )


def generate_deck(
    template_path: str,
    selected_products: List[Dict[str, Any]],
    sponsor_info: Dict[str, str]
) -> Tuple[str, str]:
    """
    Main entry point for deck generation.
    
    Generates a custom sponsorship deck by:
    1. Copying the template to a temp directory
    2. Removing non-selected product slides
    3. Adding a quote/proposal slide
    4. Converting to PDF
    
    Args:
        template_path: Path to master template.pptx
        selected_products: List of dicts with slide_index, name, price
        sponsor_info: Dict with sponsor_name, rep_name, (optional) sponsor_email
    
    Returns:
        Tuple of (pdf_path, tmpdir) - caller must clean up tmpdir after sending file
    
    Raises:
        FileNotFoundError: If template_path does not exist
        RuntimeError: If conversion to PDF fails
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")
    
    # Create temp directory
    tmpdir = tempfile.mkdtemp(prefix="humanx_deck_")
    logger.info(f"Created temp directory: {tmpdir}")
    
    try:
        # Copy template
        temp_pptx = os.path.join(tmpdir, "working.pptx")
        shutil.copy2(template_path, temp_pptx)
        logger.info(f"Copied template to: {temp_pptx}")
        
        # Load presentation
        prs = Presentation(temp_pptx)
        logger.info(f"Loaded presentation with {len(prs.slides)} slides")
        
        # Get selected slide indices
        selected_indices = {p["slide_index"] for p in selected_products}
        logger.info(f"Selected product indices: {sorted(selected_indices)}")
        
        # Collect slides to remove: all products not in selected list
        slides_to_remove = [
            idx for idx in PRODUCT_RANGE if idx not in selected_indices
        ]
        
        logger.info(f"Removing {len(slides_to_remove)} unselected product slides")
        
        # Remove slides in reverse order to avoid index shifting
        for idx in sorted(slides_to_remove, reverse=True):
            remove_slide(prs, idx)
        
        # After removing slides, recalculate the insert position
        # FRONT_MATTER_END (22) + 1 to insert after intro, but accounting for removals
        # The new position is FRONT_MATTER_END + 1
        insert_position = FRONT_MATTER_END + 1
        
        logger.info(f"Inserting quote slide at position {insert_position}")
        
        # Add quote/proposal slide
        add_quote_slide(
            prs,
            sponsor_name=sponsor_info["sponsor_name"],
            rep_name=sponsor_info["rep_name"],
            products=selected_products,
            insert_position=insert_position
        )
        
        # Save modified PPTX
        logger.info(f"Saving modified presentation to: {temp_pptx}")
        prs.save(temp_pptx)
        
        # Convert to PDF
        pdf_path = convert_to_pdf(temp_pptx, tmpdir)
        logger.info(f"Generated PDF: {pdf_path}")
        
        return pdf_path, tmpdir
    
    except Exception as e:
        logger.error(f"Error generating deck: {e}")
        # Clean up temp directory on error
        try:
            shutil.rmtree(tmpdir)
        except Exception as cleanup_error:
            logger.warning(f"Error cleaning up temp directory: {cleanup_error}")
        raise
